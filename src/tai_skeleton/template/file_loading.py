"""Content type-detection + dispatch registry for :meth:`ResourceManager.load_file`.

The type of loaded bytes is detected with Magika (a content sniff, more reliable
than an HTTP ``Content-Type`` header, which is often a generic ``octet-stream``),
then dispatched through a registry keyed by Magika label -> mime substring ->
extension. A text handler extracts the document's text; a media handler passes the
bytes through as a :data:`~tai_skeleton.template.media.MediaBlock`. Path-based
extractors (PDF, DOCX, XLSX, CSV, PPTX, ODT, EPUB) materialize the bytes to a temp
file first; the rest work on the bytes directly.

Backing dependencies are opt-in: ``pip install 'tai-skeleton[files]'``. Missing
them raises a loud install hint at import time rather than a silent skip.
"""

from __future__ import annotations

import mimetypes
import os
import tempfile
from collections.abc import Callable, Iterator
from contextlib import contextmanager
from dataclasses import dataclass
from functools import lru_cache

from tai_skeleton.template.media import MediaBlock
from tai_skeleton.template.settings import file_loading_settings

try:
    import ebooklib
    import pandas as pd
    from bs4 import BeautifulSoup
    from ebooklib import epub
    from fastmcp.utilities.types import Audio, Image
    from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader
    from magika import Magika
    from odf import teletype
    from odf.opendocument import load as odf_load
    from odf.text import P as OdfParagraph
    from pptx import Presentation
    from striprtf.striprtf import rtf_to_text
except ImportError as exc:
    raise ImportError(
        "tai-skeleton 'files' resource loading requires the 'files' optional "
        "dependencies (magika, pandas, langchain-community, pypdf, docx2txt, "
        "openpyxl, tabulate, beautifulsoup4, python-pptx, striprtf, odfpy, "
        "ebooklib). Install them with: pip install 'tai-skeleton[files]'"
    ) from exc


@lru_cache(maxsize=1)
def _magika_client() -> Magika:
    """The process-wide Magika client (constructing it loads a model, so it is
    built once and reused)."""
    return Magika()


def detect(data: bytes, source: str | None) -> tuple[str, str, str]:
    """Inspect ``data`` with Magika and return ``(label, mime, extension)``.

    ``extension`` includes the leading dot (e.g. ``".pdf"``) and is lowercased;
    it may be empty when neither Magika nor ``source`` yields one. When Magika
    cannot classify the content, ``source``'s own suffix is the extension hint.
    """
    output = _magika_client().identify_bytes(data).output
    label = output.label
    mime = output.mime_type

    ext = ""
    if output.extensions:
        ext = f".{output.extensions[0]}"
    elif label:
        ext = f".{label}"

    if source and (label in ("unknown", "undefined") or not ext):
        url_ext = os.path.splitext(source.split("?")[0])[1]
        if url_ext:
            ext = url_ext

    return label, mime, ext.lower()


def detect_mime(data: bytes) -> str | None:
    """Return the Magika-detected mime for ``data`` (``None`` when unknown)."""
    return _magika_client().identify_bytes(data).output.mime_type


@contextmanager
def _temp_file(data: bytes, suffix: str) -> Iterator[str]:
    """Materialize ``data`` to a temp file with ``suffix`` for a path-based loader,
    removing it on exit."""
    fd, path = tempfile.mkstemp(suffix=suffix)
    try:
        with os.fdopen(fd, "wb") as handle:
            handle.write(data)
        yield path
    finally:
        if os.path.exists(path):
            os.remove(path)


def _media_subtype(mime: str | None, source: str | None, kind: str) -> str:
    """The media subtype (``png``, ``wav``, ...) for a ``kind`` (``image``/``audio``)
    block, from the detected mime with a ``source``-suffix fallback."""
    if mime and mime.startswith(f"{kind}/"):
        return mime.split("/", 1)[1]
    if source:
        guessed, _ = mimetypes.guess_type(source.split("?")[0])
        if guessed and guessed.startswith(f"{kind}/"):
            return guessed.split("/", 1)[1]
    return kind


def _load_pdf(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".pdf") as path:
        return "\n".join(doc.page_content for doc in PyPDFLoader(path).load())


def _load_docx(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".docx") as path:
        return "\n".join(doc.page_content for doc in Docx2txtLoader(path).load())


def _load_xlsx(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".xlsx") as path:
        return pd.read_excel(path).to_markdown(index=False)


def _load_csv(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".csv") as path:
        return pd.read_csv(path).to_markdown(index=False)


def _load_pptx(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".pptx") as path:
        lines: list[str] = []
        for slide in Presentation(path).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)  # pyright: ignore[reportAttributeAccessIssue]
        return "\n".join(lines)


def _load_odt(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".odt") as path:
        document = odf_load(path)
        return "\n".join(teletype.extractText(p) for p in document.getElementsByType(OdfParagraph))


def _load_epub(data: bytes, mime: str | None, source: str | None) -> str:
    with _temp_file(data, ".epub") as path:
        book = epub.read_epub(path)
        parts = [
            BeautifulSoup(item.get_content(), "html.parser").get_text()
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT)
        ]
        return "\n".join(parts)


def _load_html(data: bytes, mime: str | None, source: str | None) -> str:
    return BeautifulSoup(data, "html.parser").get_text()


def _load_rtf(data: bytes, mime: str | None, source: str | None) -> str:
    return rtf_to_text(data.decode("utf-8"))


def _load_text(data: bytes, mime: str | None, source: str | None) -> str:
    # Strict decode: non-utf-8 bytes raise loudly rather than silently corrupt.
    return data.decode("utf-8")


def _load_image(data: bytes, mime: str | None, source: str | None) -> MediaBlock:
    return Image(data=data, format=_media_subtype(mime, source, "image"))


def _load_audio(data: bytes, mime: str | None, source: str | None) -> MediaBlock:
    return Audio(data=data, format=_media_subtype(mime, source, "audio"))


@dataclass(frozen=True)
class _Handler:
    """One registry entry: the Magika labels, mime substrings, and extensions it
    matches, and the loader that turns the bytes into text or a ``MediaBlock``."""

    labels: tuple[str, ...]
    mime_substrings: tuple[str, ...]
    extensions: tuple[str, ...]
    load: Callable[[bytes, str | None, str | None], str | MediaBlock]


# Order matters for the mime/extension fallback passes: an earlier entry wins.
_REGISTRY: tuple[_Handler, ...] = (
    _Handler(("pdf",), ("pdf",), (".pdf",), _load_pdf),
    _Handler(("docx",), ("wordprocessingml.document",), (".docx", ".doc"), _load_docx),
    _Handler(("xlsx",), ("spreadsheetml.sheet", "ms-excel"), (".xlsx", ".xls"), _load_xlsx),
    _Handler(("csv",), ("csv",), (".csv",), _load_csv),
    _Handler(("pptx",), ("presentationml.presentation", "ms-powerpoint"), (".pptx", ".ppt"), _load_pptx),
    _Handler(("odt",), ("opendocument.text",), (".odt",), _load_odt),
    _Handler(("epub",), ("epub",), (".epub",), _load_epub),
    _Handler(("html", "xhtml"), ("html",), (".html", ".htm"), _load_html),
    _Handler(("rtf",), ("rtf",), (".rtf",), _load_rtf),
    _Handler(
        ("txt", "markdown", "json", "xml", "yaml"),
        ("text/", "json", "xml", "yaml"),
        (".txt", ".md", ".json", ".xml", ".yaml", ".yml"),
        _load_text,
    ),
    _Handler(
        ("png", "jpeg", "gif", "webp"),
        ("image/",),
        (".png", ".jpg", ".jpeg", ".gif", ".webp"),
        _load_image,
    ),
    _Handler(("mp3", "wav"), ("audio/",), (".mp3", ".wav"), _load_audio),
)


def _dispatch(label: str, mime: str | None, ext: str) -> _Handler:
    """Resolve a handler by Magika label, then mime substring, then extension."""
    for handler in _REGISTRY:
        if label in handler.labels:
            return handler
    for handler in _REGISTRY:
        if mime and any(fragment in mime for fragment in handler.mime_substrings):
            return handler
    for handler in _REGISTRY:
        if ext and ext in handler.extensions:
            return handler
    raise ValueError(f"Unsupported file type: label={label!r} mime={mime!r} extension={ext!r}")


def load_content(data: bytes, source: str | None) -> str | MediaBlock:
    """Detect the type of ``data`` and load it into text or a ``MediaBlock``.

    ``source`` (the originating url/path, when known) supplies the suffix fallback
    for content Magika cannot classify. An unknown type raises ``ValueError``.

    The untrusted bytes are bounded by ``FILE_LOADING_MAX_BYTES`` BEFORE any
    detection or parse: a document over the cap (a storage-id source bypasses
    ``fetch_url``'s download cap) raises loudly rather than being decoded/parsed.
    """
    cap = file_loading_settings().max_bytes
    if len(data) > cap:
        raise ValueError(f"document is {len(data)} bytes, over the {cap}-byte FILE_LOADING_MAX_BYTES cap")
    label, mime, ext = detect(data, source)
    return _dispatch(label, mime, ext).load(data, mime, source)


__all__ = ["detect", "detect_mime", "load_content"]
